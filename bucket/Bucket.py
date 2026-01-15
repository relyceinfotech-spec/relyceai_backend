print("DEBUG: 1. Script started...")
import os
import sys
import csv
from dotenv import load_dotenv

# 1. Configuration
# ---------------------------------------------------------------------------
print("DEBUG: 2. Loading configuration...")
load_dotenv()

if not os.getenv("OPENROUTER_API_KEY"):
    print("Error: OPENROUTER_API_KEY not found in .env file.")
    sys.exit()

DOCS_DIR = "./documents"
DB_DIR = "./chroma_db"

# 2. Imports
# ---------------------------------------------------------------------------
print("DEBUG: 3. Importing libraries...")

# Core AI Libraries
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnableLambda
from langchain_community.retrievers import BM25Retriever
from langchain_core.documents import Document # Standard Document object

# File Handlers
from langchain_community.document_loaders import PDFPlumberLoader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
import xlrd
import mammoth
from bs4 import BeautifulSoup

print("DEBUG: 4. Imports successful!")

# ---------------------------------------------------------------------------
# CUSTOM FILE PARSERS
# ---------------------------------------------------------------------------

def parse_docx(filepath):
    """Extract text from .docx files"""
    doc = DocxDocument(filepath)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def parse_pptx(filepath):
    """Extract text from .pptx slides"""
    prs = Presentation(filepath)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def parse_excel(filepath):
    """Extract text from .xlsx files"""
    wb = load_workbook(filepath, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                text.append(" | ".join(row_text))
    return "\n".join(text)

def parse_csv(filepath):
    """Extract text from .csv files"""
    text = []
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        reader = csv.reader(f)
        for row in reader:
            text.append(" | ".join(row))
    return "\n".join(text)

def parse_html(filepath):
    """Extract clean text from .html files"""
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        soup = BeautifulSoup(f, 'html.parser')
        return soup.get_text(separator=' ', strip=True)

# ---------------------------------------------------------------------------
# UNIVERSAL LOADER FUNCTION
# ---------------------------------------------------------------------------

def load_and_process_documents():
    print(f"DEBUG: Scanning documents in {DOCS_DIR}...")
    
    if not os.path.exists(DOCS_DIR):
        os.makedirs(DOCS_DIR)
        print(f"Created {DOCS_DIR}. Please put your files there and restart.")
        sys.exit()

    raw_documents = []

    # Iterate through every file in the directory
    for root, dirs, files in os.walk(DOCS_DIR):
        for filename in files:
            filepath = os.path.join(root, filename)
            ext = filename.lower().split('.')[-1]
            
            try:
                content = ""
                # A. PDF (Using PDFPlumber)
                if ext == 'pdf':
                    loader = PDFPlumberLoader(filepath)
                    docs = loader.load()
                    raw_documents.extend(docs)
                    continue # PDFPlumber returns documents directly, so skip manual creation

                # B. Word (.docx)
                elif ext == 'docx':
                    content = parse_docx(filepath)
                
                # C. PowerPoint (.pptx)
                elif ext == 'pptx':
                    content = parse_pptx(filepath)
                
                # D. Excel (.xlsx)
                elif ext == 'xlsx':
                    content = parse_excel(filepath)
                
                # E. CSV (.csv)
                elif ext == 'csv':
                    content = parse_csv(filepath)
                
                # F. HTML (.html / .htm)
                elif ext in ['html', 'htm']:
                    content = parse_html(filepath)
                
                # G. Text (.txt, .md, .py)
                elif ext in ['txt', 'md', 'py']:
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()

                # If content was extracted, wrap it in a Document object
                if content:
                    doc = Document(page_content=content, metadata={"source": filename})
                    raw_documents.append(doc)
                    print(f"  - Loaded: {filename}")
            
            except Exception as e:
                print(f"  ! Skipped {filename}: {e}")

    if not raw_documents:
        print("No documents successfully loaded. Please check your files.")
        sys.exit()

    print(f"Successfully loaded {len(raw_documents)} documents. Splitting text...")
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200, add_start_index=True)
    splits = text_splitter.split_documents(raw_documents)
    return splits

# ---------------------------------------------------------------------------
# RAG COMPONENTS
# ---------------------------------------------------------------------------

def setup_vectorstore(splits):
    print("Initializing Embeddings (CPU mode)...")
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

    print("Creating/Updating Vector Database...")
    vectorstore = Chroma.from_documents(documents=splits, embedding=embeddings, persist_directory=DB_DIR)
    return vectorstore

def hybrid_retrieval_func(query, bm25_retriever, chroma_retriever):
    """Combines Keyword (BM25) and Semantic (Chroma) results."""
    bm25_docs = bm25_retriever.invoke(query)
    chroma_docs = chroma_retriever.invoke(query)
    
    combined_docs = []
    seen_content = set()
    
    max_len = max(len(bm25_docs), len(chroma_docs))
    for i in range(max_len):
        if i < len(bm25_docs):
            doc = bm25_docs[i]
            if doc.page_content not in seen_content:
                combined_docs.append(doc)
                seen_content.add(doc.page_content)
        if i < len(chroma_docs):
            doc = chroma_docs[i]
            if doc.page_content not in seen_content:
                combined_docs.append(doc)
                seen_content.add(doc.page_content)
                
    return combined_docs[:6]

def create_advanced_rag_chain(vectorstore, splits):
    print("Setting up LLM and Chains...")
    llm = ChatOpenAI(
        model_name="openai/gpt-4o-mini", 
        
        # Configuration for OpenRouter
        openai_api_key=os.getenv("OPENROUTER_API_KEY"),
        openai_api_base="https://openrouter.ai/api/v1",
        
        # Standard Settings
        temperature=0,
    )

    bm25_retriever = BM25Retriever.from_documents(splits)
    bm25_retriever.k = 3
    chroma_retriever = vectorstore.as_retriever(search_kwargs={"k": 3})

    rewrite_prompt = ChatPromptTemplate.from_template(
        "Refine this query for better search results. Query: {x} \nOptimized Query:"
    )
    rewriter_chain = rewrite_prompt | llm | StrOutputParser()

    # UPDATED SYSTEM PROMPT WITH RELYCE AI IDENTITY
    # Answer Chain
    system_prompt = """You are the friendly and intelligent AI assistant for **Relyce AI**. Your goal is to be helpful, clear, and accurate.

    **Who You Are (Relyce AI):**
    You are a smart AI platform designed to make work easier. You use artificial intelligence to automate difficult tasks, analyze data precisely, and give useful advice. Your mission is to help people work better and make smarter decisions using powerful, simple-to-use AI.

    **Your Knowledge & Capabilities:**
    Although you speak in simple English, you are very knowledgeable in these areas:
    1. **Technology:** You understand software, data, and computer security.
    2. **Ethics:** You care about using AI safely and protecting privacy.
    3. **Rules & Laws:** You know about general legal rules (but remember, you are not a lawyer and cannot give legal advice).

    **Rules for Answering:**
    * **Be Simple & Clear:** Use plain English. Avoid complicated jargon unless necessary.
    * **Stick to the Facts:** Answer **only** using the information provided in the "Context Documents" below. Do not guess or use outside knowledge.
    * **Be Honest:** If the answer is not in the documents, just say: "I'm sorry, but I cannot find that information in the documents provided."
    * **Be Precise:** When you mention facts or numbers, use the exact details from the documents.
    * **Friendly Tone:** Be polite and helpful, like a smart assistant speaking to a normal person.

    Context Documents:
    {context}"""

    answer_prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        ("human", "{input}"),
    ])

    def format_docs(docs):
        return "\n\n".join(doc.page_content for doc in docs)

    hybrid_retriever_runnable = RunnableLambda(
        lambda q: hybrid_retrieval_func(q, bm25_retriever, chroma_retriever)
    )

    rag_chain = (
        {
            "context": (lambda x: x["input"]) | rewriter_chain | hybrid_retriever_runnable | format_docs,
            "input": lambda x: x["input"]
        }
        | answer_prompt 
        | llm 
        | StrOutputParser()
    )
    return rag_chain

# 3. Execution Loop
# ---------------------------------------------------------------------------
print("DEBUG: 5. Entering main execution...")

if __name__ == "__main__":
    print("\n--- STARTING RELYCE AI UNIVERSAL ENGINE ---\n")
    
    splits = load_and_process_documents()
    
    if os.path.exists(DB_DIR):
        print("Loading existing Chroma database...")
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        vectorstore = Chroma(persist_directory=DB_DIR, embedding_function=embeddings)
    else:
        vectorstore = setup_vectorstore(splits)

    chain = create_advanced_rag_chain(vectorstore, splits)

    print("\n" + "="*50)
    print("   RELYCE AI: READY")
    print("="*50)
    print("Type 'exit' to quit.\n")

    while True:
        try:
            query = input("User Query: ")
            if query.lower() in ["exit", "quit"]:
                break
            
            print("\nOptimizing query & Consulting knowledge base...\n")
            response = chain.invoke({"input": query})
            print("Response:")
            print(response)
            print("-" * 50)
        except KeyboardInterrupt:
            print("\nExiting...")
            break
        except Exception as e:
            print(f"Error: {e}")