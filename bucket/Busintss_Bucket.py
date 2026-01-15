print("DEBUG: 1. Script started...")
import os
import sys
import csv
from dotenv import load_dotenv

# 1. Configuration
# ---------------------------------------------------------------------------
print("DEBUG: 2. Loading configuration...")
load_dotenv()

if not os.getenv("OPENROUTER_API_KEY"):
    print("Warning: OPENROUTER_API_KEY not found in .env file. RAG features may not work.")

DOCS_DIR = "./documents" 
DB_DIR = "./chroma_db"

# 2. Imports
# ---------------------------------------------------------------------------
print("DEBUG: 3. Importing libraries...")

# Core AI Libraries
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnableLambda
from langchain_community.retrievers import BM25Retriever
from langchain_core.documents import Document # Standard Document object

# File Handlers
from langchain_community.document_loaders import PDFPlumberLoader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
# xlrd and mammoth removed - not needed
from bs4 import BeautifulSoup

print("DEBUG: 4. Imports successful!")

# ---------------------------------------------------------------------------
# CUSTOM FILE PARSERS
# ---------------------------------------------------------------------------

def parse_docx(filepath):
    """Extract text from .docx files"""
    doc = DocxDocument(filepath)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def parse_pptx(filepath):
    """Extract text from .pptx slides"""
    prs = Presentation(filepath)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def parse_excel(filepath):
    """Extract text from .xlsx files"""
    wb = load_workbook(filepath, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                text.append(" | ".join(row_text))
    return "\n".join(text)

def parse_csv(filepath):
    """Extract text from .csv files"""
    text = []
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        reader = csv.reader(f)
        for row in reader:
            text.append(" | ".join(row))
    return "\n".join(text)

def parse_html(filepath):
    """Extract clean text from .html files"""
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        soup = BeautifulSoup(f, 'html.parser')
        return soup.get_text(separator=' ', strip=True)

# ---------------------------------------------------------------------------
# UNIVERSAL LOADER FUNCTION
# ---------------------------------------------------------------------------

def load_and_process_documents():
    print(f"DEBUG: Scanning documents in {DOCS_DIR}...")
    
    if not os.path.exists(DOCS_DIR):
        os.makedirs(DOCS_DIR)
        print(f"Created {DOCS_DIR}. No documents to load.")
        return []  # Return empty list instead of crashing

    raw_documents = []

    # Iterate through every file in the directory
    for root, dirs, files in os.walk(DOCS_DIR):
        for filename in files:
            filepath = os.path.join(root, filename)
            ext = filename.lower().split('.')[-1]
            
            try:
                content = ""
                # A. PDF (Using PDFPlumber)
                if ext == 'pdf':
                    loader = PDFPlumberLoader(filepath)
                    docs = loader.load()
                    raw_documents.extend(docs)
                    continue # PDFPlumber returns documents directly, so skip manual creation

                # B. Word (.docx)
                elif ext == 'docx':
                    content = parse_docx(filepath)
                
                # C. PowerPoint (.pptx)
                elif ext == 'pptx':
                    content = parse_pptx(filepath)
                
                # D. Excel (.xlsx)
                elif ext == 'xlsx':
                    content = parse_excel(filepath)
                
                # E. CSV (.csv)
                elif ext == 'csv':
                    content = parse_csv(filepath)
                
                # F. HTML (.html / .htm)
                elif ext in ['html', 'htm']:
                    content = parse_html(filepath)
                
                # G. Text (.txt, .md, .py)
                elif ext in ['txt', 'md', 'py']:
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()

                # If content was extracted, wrap it in a Document object
                if content:
                    doc = Document(page_content=content, metadata={"source": filename})
                    raw_documents.append(doc)
                    print(f"  - Loaded: {filename}")
            
            except Exception as e:
                print(f"  ! Skipped {filename}: {e}")

    if not raw_documents:
        print("No documents found. Returning empty list.")
        return []  # Return empty list instead of crashing

    print(f"Successfully loaded {len(raw_documents)} documents. Splitting text...")
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200, add_start_index=True)
    splits = text_splitter.split_documents(raw_documents)
    return splits

# ---------------------------------------------------------------------------
# RAG COMPONENTS
# ---------------------------------------------------------------------------

def setup_vectorstore(splits):
    print("Initializing Embeddings (CPU mode)...")
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

    print("Creating/Updating Vector Database...")
    vectorstore = Chroma.from_documents(documents=splits, embedding=embeddings, persist_directory=DB_DIR)
    return vectorstore

def hybrid_retrieval_func(query, bm25_retriever, chroma_retriever):
    """Combines Keyword (BM25) and Semantic (Chroma) results."""
    bm25_docs = bm25_retriever.invoke(query)
    chroma_docs = chroma_retriever.invoke(query)
    
    combined_docs = []
    seen_content = set()
    
    max_len = max(len(bm25_docs), len(chroma_docs))
    for i in range(max_len):
        if i < len(bm25_docs):
            doc = bm25_docs[i]
            if doc.page_content not in seen_content:
                combined_docs.append(doc)
                seen_content.add(doc.page_content)
        if i < len(chroma_docs):
            doc = chroma_docs[i]
            if doc.page_content not in seen_content:
                combined_docs.append(doc)
                seen_content.add(doc.page_content)
                
    return combined_docs[:6]

def create_advanced_rag_chain(vectorstore, splits):
    print("Setting up LLM and Chains...")
    llm = ChatOpenAI(
        model_name="gpt-5-nano", 
        openai_api_key=os.getenv("OPENROUTER_API_KEY"),
        openai_api_base="https://openrouter.ai/api/v1",
        temperature=0,
    )

    bm25_retriever = BM25Retriever.from_documents(splits)
    bm25_retriever.k = 3
    chroma_retriever = vectorstore.as_retriever(search_kwargs={"k": 3})

    rewrite_prompt = ChatPromptTemplate.from_template(
        "Refine this query for better search results. Query: {x} \nOptimized Query:"
    )
    rewriter_chain = rewrite_prompt | llm | StrOutputParser()

    # UPDATED SYSTEM PROMPT WITH RELYCE AI IDENTITY
    system_prompt = """You are a highly accomplished and multi-faceted AI assistant, functioning as an **elite consultant and strategic advisor** for businesses and startups. Your persona embodies the collective expertise of a Chief Operating Officer, a Head of Legal, a Chief Technology Officer, and a Chief Ethics Officer.

    **Core Mandate:**
    You must provide zero-hallucination, fact-based guidance operating with:

    1. **Business Acumen:** Deep understanding of market dynamics, growth strategies, operational efficiencies, and financial models.
    2. **Startup Savvy:** Intimate knowledge of the startup lifecycle, fundraising, lean methodologies, and scaling challenges.
    3. **Technical Proficiency:** Ability to discuss technology stacks, software development, data analytics, and cybersecurity with precision.
    4. **Ethical Integrity:** A commitment to responsible AI usage, data privacy, and understanding the societal impact of business decisions.
    5. **Legal Prudence:** Awareness of legal frameworks, IP, and compliance. (Note: You are not a lawyer; identify considerations but do not provide legal advice.)
    6. **Corporate Identity (Relyce AI):** You are the proprietary AI engine of **Relyce AI**, a state-of-the-art platform engineered to revolutionize enterprise workflows through intelligent automation. Your purpose is to automate complex operational tasks, execute high-precision data analysis, and generate actionable strategic insights. You are dedicated to enhancing organizational productivity and driving smarter, data-backed business decisions through a powerful and intuitive AI solution.

    **Strict Guidelines for Response Generation:**
    * **Greetings:** You may answer simple greetings (like 'hi', 'hello') politely and professionally.
    * **Context-Bound:** For all other queries, your answers must be derived **solely and exclusively** from the provided retrieved context. Do not infer, speculate, or use external knowledge.
    * **Zero Hallucination:** If the provided context does not contain sufficient information to answer the question, state clearly: "Based on the available documents, the information to fully address this specific query is not present."
    * **Conciseness & Precision:** Be direct, highly precise, and professional. Avoid filler.
    * **Detail-Oriented:** Provide specific details, figures, and sources (e.g., "Document X, Page Y") when the context supports it.
    * **Synthesis:** Synthesize multiple relevant pieces of context into a strategic response actionable for a business leader.
    * **Tone:** Maintain a professional, authoritative, and advisory tone.

    Context Documents:
    {context}"""

    answer_prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        ("human", "{input}"),
    ])

    def format_docs(docs):
        return "\n\n".join(doc.page_content for doc in docs)

    hybrid_retriever_runnable = RunnableLambda(
        lambda q: hybrid_retrieval_func(q, bm25_retriever, chroma_retriever)
    )

    rag_chain = (
        {
            "context": (lambda x: x["input"]) | rewriter_chain | hybrid_retriever_runnable | format_docs,
            "input": lambda x: x["input"]
        }
        | answer_prompt 
        | llm 
        | StrOutputParser()
    )
    return rag_chain

# 3. Execution Loop
# ---------------------------------------------------------------------------
print("DEBUG: 5. Entering main execution...")

if __name__ == "__main__":
    print("\n--- STARTING RELYCE AI UNIVERSAL ENGINE ---\n")
    
    splits = load_and_process_documents()
    
    if os.path.exists(DB_DIR):
        print("Loading existing Chroma database...")
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        vectorstore = Chroma(persist_directory=DB_DIR, embedding_function=embeddings)
    else:
        vectorstore = setup_vectorstore(splits)

    chain = create_advanced_rag_chain(vectorstore, splits)

    print("\n" + "="*50)
    print("   RELYCE AI: READY")
    print("="*50)
    print("Type 'exit' to quit.\n")

    while True:
        try:
            query = input("User Query: ")
            if query.lower() in ["exit", "quit"]:
                break
            
            print("\nOptimizing query & Consulting knowledge base...\n")
            response = chain.invoke({"input": query})
            print("Response:")
            print(response)
            print("-" * 50)
        except KeyboardInterrupt:
            print("\nExiting...")
            break
        except Exception as e:
            print(f"Error: {e}")